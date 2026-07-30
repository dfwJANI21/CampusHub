from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Apply a dark theme-like approach by customizing slide backgrounds if possible, 
# but for simplicity, we will use default layouts and add content.
# Layouts: 0-Title, 1-Title & Content, 2-Section Header, etc.

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Campus Event & Club Hub"
subtitle.text = "A Centralized Platform for College Clubs and Events\n\nPresented by:\nChavda Pruthviraj, Jainam Jani, Pansare Dhruv, Gaurav Sengar"

# Slide 2: Problem Statement
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "The Problem"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Campus life can be disconnected:"
p = tf.add_paragraph()
p.text = "Students often miss out on events because information is scattered."
p.level = 1
p = tf.add_paragraph()
p.text = "Clubs struggle to reach the student body effectively."
p.level = 1
p = tf.add_paragraph()
p.text = "Reliance on chaotic WhatsApp groups and easily ignored physical bulletin boards."
p.level = 1

# Slide 3: Our Solution
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Our Solution"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Introducing the Campus Event & Club Hub:"
p = tf.add_paragraph()
p.text = "A modern, unified web application dedicated to campus activities."
p.level = 1
p = tf.add_paragraph()
p.text = "Bridges the communication gap between clubs and students."
p.level = 1
p = tf.add_paragraph()
p.text = "Provides a single source of truth for what's happening on campus."
p.level = 1

# Slide 4: Key Features for Students
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Key Features (For Students)"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Enhancing student engagement:"
p = tf.add_paragraph()
p.text = "Discover upcoming events easily through a beautiful feed."
p.level = 1
p = tf.add_paragraph()
p.text = "Get detailed information: Time, Location, and Descriptions."
p.level = 1
p = tf.add_paragraph()
p.text = "Explore the active clubs on campus."
p.level = 1

# Slide 5: Key Features for Clubs
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Key Features (For Clubs)"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Empowering student organizations:"
p = tf.add_paragraph()
p.text = "Easy-to-use 'Post Event' form to announce activities instantly."
p.level = 1
p = tf.add_paragraph()
p.text = "Direct outreach to a broader student audience."
p.level = 1
p = tf.add_paragraph()
p.text = "Dedicated branding and visibility within the hub."
p.level = 1

# Slide 6: Technology Stack
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Technology Stack"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Built with modern and robust tools:"
p = tf.add_paragraph()
p.text = "Backend: Python (Flask Framework)"
p.level = 1
p = tf.add_paragraph()
p.text = "Database: SQLite with SQLAlchemy ORM"
p.level = 1
p = tf.add_paragraph()
p.text = "Frontend: HTML5, CSS3, Vanilla JavaScript"
p.level = 1
p = tf.add_paragraph()
p.text = "Design: Custom Dark Theme with Glassmorphism UI"
p.level = 1

# Slide 7: System Architecture
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "System Architecture"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "MVC (Model-View-Controller) Pattern:"
p = tf.add_paragraph()
p.text = "Model: SQLAlchemy handles database schemas and queries."
p.level = 1
p = tf.add_paragraph()
p.text = "View: Jinja2 templates render dynamic HTML pages."
p.level = 1
p = tf.add_paragraph()
p.text = "Controller: Flask routes process requests and handle business logic."
p.level = 1

# Slide 8: Database Schema
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Database Schema"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Relational structure:"
p = tf.add_paragraph()
p.text = "Club Model: ID, Name, Description, Logo URL"
p.level = 1
p = tf.add_paragraph()
p.text = "Event Model: ID, Title, Date, Time, Location, Description, Club ID (Foreign Key)"
p.level = 1
p = tf.add_paragraph()
p.text = "Relationship: One-to-Many (One Club can host Many Events)"
p.level = 1

# Slide 9: Future Enhancements
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Future Enhancements"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Looking ahead:"
p = tf.add_paragraph()
p.text = "User Authentication (Student logins & Club Admin dashboards)."
p.level = 1
p = tf.add_paragraph()
p.text = "RSVP and Event Ticketing System."
p.level = 1
p = tf.add_paragraph()
p.text = "Email and Push Notifications for event reminders."
p.level = 1

# Slide 10: Conclusion & Q&A
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Enhancing campus life through better connectivity."
p = tf.add_paragraph()
p.text = "Thank You!"
p.level = 1
p = tf.add_paragraph()
p.text = "Questions?"
p.level = 1

# Save presentation
prs.save('Campus_Hub_Presentation.pptx')
print("Presentation successfully created at Campus_Hub_Presentation.pptx")
